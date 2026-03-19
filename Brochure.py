import streamlit as st
from pptx import Presentation
import requests
from copy import deepcopy
import io

st.set_page_config(page_title="Gestionnaire de Brochure", layout="wide")

if 'rows' not in st.session_state:
    st.session_state.rows = []

TOKEN = st.secrets["FORMATION_TOKEN"]
API_URL = st.secrets["URL_FORMATION"]


headers = {
    "Authorization": f"Bearer {TOKEN}",
    "Content-Type": "application/json"
}

def chunk_list(lst, size=10):
    for i in range(0, len(lst), size):
        yield lst[i:i+size]

def duplicate_slide(prs, slide_index):
    source_slide = prs.slides[slide_index]

    # Utilise le même layout que la slide source
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Supprimer les placeholders par défaut
    for shape in list(new_slide.shapes):
        el = shape.element
        el.getparent().remove(el)

    # Copier toutes les shapes
    for shape in source_slide.shapes:
        new_shape = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            new_shape, 'p:extLst'
        )

    return new_slide

def move_slide_to(prs, slide, target_index):
    slides = prs.slides._sldIdLst

    # retrouver l'élément XML de la slide
    for i, sldId in enumerate(slides):
        if prs.slides[i] == slide:
            slide_id = sldId
            current_index = i
            break
    else:
        return  # slide non trouvée

    # déplacer
    slides.remove(slide_id)
    slides.insert(target_index, slide_id)

def update_text_preserve_style(shape, new_text):
    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        return
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = str(new_text) if new_text is not None else ""
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = str(new_text) if new_text is not None else ""

def delete_slide(prs, slide_index):
    slide_id = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(slide_id)

@st.dialog("Ajouter une nouvelle formation !")
def add_formation():
    col1, col2 = st.columns(2)
    with col1:
        add_nom = st.text_input("Nom de la formation")
        add_type = st.selectbox("Type ", ["BTS", "BACHELOR", "MASTERE", "BAC+6", "DOCTORATE"])
        add_langue = st.text_input("Langue(s) de la formation")
        add_langue_p = st.selectbox("Langue de la page de la brochure ", ["Français", "English"])
    with col2:
        add_desc = st.text_area("Description ", height=100)
        add_stage = st.text_input("Stage ")
        add_admission = st.text_input("Admission ")

    with st.expander("Détails supplémentaires "):
        c1, c2, c3 = st.columns(3)
        af1 = c1.text_area("Point Fort 1 ")
        af2 = c1.text_area("Point Fort 2 ")
        af3 = c1.text_area("Point Fort 3 ")
        ae1 = c2.text_area("Enseignement 1 ")
        ae2 = c2.text_area("Enseignement 2 ")
        ae3 = c2.text_area("Enseignement 3 ")
        amet = c3.text_area("Métiers ")

    if st.button("Enregistrer la formation", type="primary"):
        new_row = {
            "Nom": add_nom,
            "Type": add_type,
            "Langues": add_langue,
            "Langue_Formation": add_langue_p,
            "Description": add_desc,
            "Stage": add_stage,
            "Admission": add_admission,
            "PointFort1": af1, "PointFort2": af2, "PointFort3": af3,
            "Enseignement1": ae1, "Enseignement2": ae2, "Enseignement3": ae3,
            "Metier": amet,
            "id": None
        }

        missing = [name for name, value in new_row.items() if name != "id" and (not value or not str(value).strip())]
        if missing:
            st.error(f"Veuillez remplir les champs suivants : {', '.join(missing)}")
        else:
            st.session_state.rows.append(new_row)
            st.success("Formation ajoutée au cache !")
            st.rerun()

@st.cache_data(ttl=600)
def fetch_data():
    resp = requests.get(f"{API_URL}/data", headers=headers, params={"maxRecords": 100, "view": "Grid view"})
    if resp.status_code == 200:
        data = resp.json()
        return [r for r in data], [r["id"] for r in data]
    print(resp.status_code)
    print(f"{API_URL}/data")
    return [], []

if not st.session_state.rows:
    rows, ids = fetch_data()
    st.session_state.rows = rows
    st.session_state.ids = ids

############### Interface Streamlit #################

st.title("Gestionnaire de Brochure Formations")

tab1, tab2, tab3 = st.tabs(["Édition des Formations", "Vue d'ensemble", "Génération"])

# --- TAB 1 ---
with tab1:
    st.subheader("Ajouter une formation")
    if st.button("Ajouter une formation", type="primary"):
        add_formation()
    
    st.subheader("Modifier une formation")
    formation_names = [r.get("Nom", "Sans nom") for r in st.session_state.rows]
    selected_name = st.selectbox("Choisir une formation à modifier", formation_names)

    idx = formation_names.index(selected_name)
    current_row = st.session_state.rows[idx]

    if st.button("Supprimer la formation", use_container_width=True):
        st.session_state.rows.pop(idx)
        st.warning(f"La formation {selected_name} a été retirée du cache.")
        st.rerun()

    col1, col2 = st.columns(2)
    with col1:
        new_nom = st.text_input("Nom de la formation", current_row.get("Nom", ""))
        new_type = st.selectbox("Type", ["BTS", "BACHELOR", "MASTERE", "BAC+6", "DOCTORATE"], 
                               index=["BTS", "BACHELOR", "MASTERE", "BAC+6", "DOCTORATE"].index(current_row.get("Type", "BTS")))
        new_langue = st.text_input("Langue(s) de la formation", current_row.get("Langues", ""))
        new_langue_p = st.selectbox("Langue de la page de la brochure", ["Français", "English"], 
                                   index=0 if current_row.get("Langue_Formation") == "Français" else 1)
    with col2:
        new_desc = st.text_area("Description", current_row.get("Description", ""), height=100)
        new_stage = st.text_input("Stage", current_row.get("Stage", ""))
        new_admission = st.text_input("Admission", current_row.get("Admission", ""))

    with st.expander("Détails supplémentaires"):
        c1, c2, c3 = st.columns(3, border=True)
        f1 = c1.text_area("Point Fort 1", current_row.get("PointFort1", ""))
        f2 = c1.text_area("Point Fort 2", current_row.get("PointFort2", ""))
        f3 = c1.text_area("Point Fort 3", current_row.get("PointFort3", ""))
        e1 = c2.text_area("Enseignement 1", current_row.get("Enseignement1", ""))
        e2 = c2.text_area("Enseignement 2", current_row.get("Enseignement2", ""))
        e3 = c2.text_area("Enseignement 3", current_row.get("Enseignement3", ""))
        met = c3.text_area("Métiers", current_row.get("Metier", ""))

    if st.button("Enregistrer les modifications", type="primary"):
        st.session_state.rows[idx].update({
            "Nom": new_nom, "Type": new_type, "Langues": new_langue, "Langue_Formation": new_langue_p, "Stage": new_stage, "Admission": new_admission,
            "Description": new_desc, "PointFort1": f1, "PointFort2": f2, "PointFort3": f3, "Enseignement1": e1, "Enseignement2": e2, "Enseignement3": e3, "Metier": met
        })
        st.success(f"Modifications pour {new_nom} enregistrées temporairement.")


# --- TAB 2 ---
with tab2:
    st.subheader("Tableau des données")
    st.dataframe(st.session_state.rows)

    if st.button("Envoyer les données", type="primary"):
        with st.spinner("Mise à jour en cours..."):
            requests.delete(f"{API_URL}/data", headers=headers)
            requests.delete(url, headers=headers)

            send = requests.post(f"{API_URL}/data", headers=headers, json=st.session_state.rows)
            send = requests.post(url, headers=headers, json=st.session_state.rows)

            if send.status_code == 200:
                st.success("Synchronisation terminée !")
                st.cache_data.clear() # IMPORTANT : On vide le cache car les données ont changé
                
                # On force la mise à jour de la session pour que l'onglet 1 voie les changements
                rows, ids = fetch_data() 
                st.session_state.rows = rows
                st.session_state.ids = ids
                st.rerun() 
            else:
                st.error(f"Erreur : {send.text}")

# --- TAB 3 ---
with tab3:
    st.info("Vérifiez que vos modifications sont bien dans le tableau (Onglet 2) avant de générer.")
    
    if st.button("Générer la Brochure (.pptx)", type="primary"):
        with st.spinner("Création de la brochure..."):
            prs = Presentation("template.pptx")
            created_slides = []
            
            # --- CREATION DES SLIDES ---
            for row in st.session_state.rows:
                mapping = {
                    "TextBox 48": row.get("Nom"), "TextBox 50": row.get("Type"),
                    "TextBox 34": row.get("Langues"), "TextBox 38": row.get("Stage"),
                    "TextBox 42": row.get("Description"), "TextBox 33": row.get("PointFort1"),
                    "TextBox 37": row.get("PointFort2"), "TextBox 40": row.get("PointFort3"),
                    "TextBox 32": row.get("Enseignement1"), "TextBox 36": row.get("Enseignement2"),
                    "TextBox 41": row.get("Enseignement3"), "TextBox 35": row.get("Metier"),
                    "TextBox 39": row.get("Admission"),
                }
                
                # Sélection du modèle selon le Type
                model_idx = {"BTS": 10, "BACHELOR": 12, "MASTERE": 15, "BAC+6": 17, "DOCTORATE": 19}.get(row["Type"], 12)
                
                slide = duplicate_slide(prs, model_idx)
                created_slides.append((slide, model_idx))
                
                for shape in slide.shapes:
                    if shape.name in mapping:
                        update_text_preserve_style(shape, mapping[shape.name])

            # Organisation et suppression des modèles
            by_model = {10: [], 12: [], 15: [], 17: [], 19: []}
            for slide, m_idx in created_slides:
                by_model[m_idx].append(slide)

            # Supprimer les originaux
            for index in sorted([10, 12, 15, 17, 19], reverse=True):
                delete_slide(prs, index)

            # Replacer les nouvelles slides par blocs
            # Index après suppression : BTS=10, BACH=11, MAST=13, B6=14, DOC=15
            order = [(19, 15), (17, 14), (15, 13), (12, 11), (10, 10)]
            for original_idx, target_idx in order:
                for s in reversed(by_model[original_idx]):
                    move_slide_to(prs, s, target_idx)

            # Recensement pour le sommaire
            pages_map = {}
            for i, sld in enumerate(prs.slides):
                for shape in sld.shapes:
                    if shape.name == "TextBox 48":
                        name = shape.text_frame.text.strip()
                        if name: pages_map[name] = str(i + 1)

            # Remplissage du sommaire
            sommaire_slide = prs.slides[5]
            s_noms = {k: [] for k in ["BTS","BACH_FR","BACH_EN","MAST_FR","MAST_EN","B6_FR","B6_EN","DOC"]}
            s_pgs = {k: [] for k in ["BTS_P","BACH_FR_P","BACH_EN_P","MAST_FR_P","MAST_EN_P","B6_FR_P","B6_EN_P","DOC_P"]}

            for row in st.session_state.rows:
                n, t, l = row.get("Nom"), row.get("Type"), row.get("Langue_Formation")
                p = pages_map.get(n, "-")
                
                if t == "BTS": k = "BTS"
                elif t == "BACHELOR": k = "BACH_EN" if l == "English" else "BACH_FR"
                elif t == "MASTERE": k = "MAST_EN" if l == "English" else "MAST_FR"
                elif t == "BAC+6": k = "B6_EN" if l == "English" else "B6_FR"
                elif t == "DOCTORATE": k = "DOC"
                else: continue
                
                s_noms[k].append(n)
                s_pgs[k+"_P"].append("p."+p)

            # Mapping final des TextBox (Noms + Pages)
            mapping_sommaire = {
                "TextBox 8": "\n".join(s_noms["BTS"] + s_noms["BACH_FR"]),
                "TextBox 21": "\n".join(s_pgs["BTS_P"] + s_pgs["BACH_FR_P"]),
                "TextBox 5": "\n".join(s_noms["BACH_EN"]),
                "TextBox 24": "\n".join(s_pgs["BACH_EN_P"]),
                "TextBox 9": "\n".join(s_noms["MAST_FR"] + s_noms["B6_FR"]),
                "TextBox 22": "\n".join(s_pgs["MAST_FR_P"] + s_pgs["B6_FR_P"]),
                "TextBox 6": "\n".join(s_noms["MAST_EN"] + s_noms["B6_EN"]),
                "TextBox 19": "\n".join(s_pgs["MAST_EN_P"] + s_pgs["B6_EN_P"]),
                "TextBox 7": "\n".join(s_noms["DOC"]),
                "TextBox 20": "\n".join(s_pgs["DOC_P"]),
                "TextBox 10": "\n".join(s_noms["DOC"]),
                "TextBox 23": "\n".join(s_pgs["DOC_P"]),
            }

            for shape in sommaire_slide.shapes:
                if shape.name in mapping_sommaire:
                    update_text_preserve_style(shape, mapping_sommaire[shape.name])

            # Pagination
            for i, sld in enumerate(prs.slides):
                for shape in sld.shapes:
                    if shape.name == "TextBox 99":
                        update_text_preserve_style(shape, i + 1)
            
            # Sauvegarde en mémoire pour Streamlit
            binary_output = io.BytesIO()
            prs.save(binary_output)
            
            st.success("Brochure prête !")
            st.download_button(
                label="Télécharger la Brochure",
                data=binary_output.getvalue(),
                file_name="Brochure_Formations.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

with st.sidebar:
    st.pdf("template.pdf")